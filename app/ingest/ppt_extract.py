"""PowerPoint ingest: slide JPEG export (LibreOffice) + text/notes (python-pptx)."""

from __future__ import annotations

import logging
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger("ppt-extract")

_MAX_SLIDE_TEXT_CHARS = 8000
_JPEG_QUALITY = 88
_JPEG_MAX_WIDTH = 1920


@dataclass(frozen=True)
class PptSlideData:
    slide_index: int  # 0-based
    slide_num: int  # 1-based (for pages[] compatibility)
    source_text: str
    image_path: Path  # absolute path to slide_XXX.jpg


def _find_soffice() -> str:
    env = (os.getenv("SOFFICE_PATH") or "").strip()
    if env and Path(env).is_file():
        return env
    for name in ("soffice", "soffice.exe"):
        found = shutil.which(name)
        if found:
            return found
    # Common Windows LibreOffice installs
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if Path(candidate).is_file():
            return candidate
    raise FileNotFoundError(
        "LibreOffice (soffice) not found. Install LibreOffice or set SOFFICE_PATH "
        "to the soffice executable."
    )


def _convert_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    """Convert .pptx to PDF via LibreOffice headless."""
    soffice = _find_soffice()
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--norestore",
        "--convert-to",
        "pdf",
        "--outdir",
        str(out_dir),
        str(pptx_path),
    ]
    logger.info("LibreOffice convert: %s", " ".join(cmd))
    try:
        proc = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=300,
            check=False,
        )
    except subprocess.TimeoutError as exc:
        raise RuntimeError("LibreOffice conversion timed out after 300s") from exc
    if proc.returncode != 0:
        err = (proc.stderr or proc.stdout or "").strip()
        raise RuntimeError(f"LibreOffice conversion failed: {err or proc.returncode}")

    pdf_name = pptx_path.with_suffix(".pdf").name
    pdf_path = out_dir / pdf_name
    if not pdf_path.is_file():
        # LO may use stem only
        candidates = list(out_dir.glob("*.pdf"))
        if len(candidates) == 1:
            pdf_path = candidates[0]
        else:
            raise RuntimeError(
                f"LibreOffice did not produce expected PDF in {out_dir}"
            )
    return pdf_path


def _render_pdf_to_jpegs(pdf_path: Path, slides_dir: Path) -> list[Path]:
    """Render each PDF page to slide_NNN.jpg."""
    slides_dir.mkdir(parents=True, exist_ok=True)
    # LibreOffice PDFs often trigger harmless MuPDF structure-tree warnings.
    if hasattr(fitz, "TOOLS"):
        fitz.TOOLS.mupdf_display_errors(False)
    doc = fitz.open(str(pdf_path))
    paths: list[Path] = []
    try:
        for i in range(doc.page_count):
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=220, alpha=False)
            if pix.width > _JPEG_MAX_WIDTH:
                scale = _JPEG_MAX_WIDTH / pix.width
                mat = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=mat, alpha=False)
            out = slides_dir / f"slide_{i + 1:03d}.jpg"
            pix.save(str(out), jpg_quality=_JPEG_QUALITY)
            paths.append(out)
    finally:
        doc.close()
    return paths


def _shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return (shape.text or "").strip()


def _slide_visible_text(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        t = _shape_text(shape)
        if t:
            parts.append(t)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for sub in shape.shapes:
                    st = _shape_text(sub)
                    if st:
                        parts.append(st)
            except Exception:
                pass
    return "\n".join(parts).strip()


def _extract_slide_texts(pptx_path: Path) -> list[str]:
    prs = Presentation(str(pptx_path))
    texts: list[str] = []
    for slide in prs.slides:
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        visible = _slide_visible_text(slide)
        chunks = []
        if notes:
            chunks.append(f"[Speaker notes]\n{notes}")
        if visible:
            chunks.append(f"[On-slide text]\n{visible}")
        combined = "\n\n".join(chunks).strip()
        if len(combined) > _MAX_SLIDE_TEXT_CHARS:
            combined = combined[:_MAX_SLIDE_TEXT_CHARS] + "…"
        texts.append(combined or f"(Slide {len(texts) + 1} — visual content; refer to the slide image.)")
    return texts


def ingest_pptx(
    pptx_path: Path,
    *,
    slides_dir: Path,
    session_id: str,
) -> list[PptSlideData]:
    """
    Export slide JPEGs and per-slide text for a tutoring session.

    Raises FileNotFoundError / RuntimeError on failure.
    """
    pptx_path = Path(pptx_path).resolve()
    slides_dir = Path(slides_dir).resolve()
    slides_dir.mkdir(parents=True, exist_ok=True)

    slide_texts = _extract_slide_texts(pptx_path)

    with tempfile.TemporaryDirectory(prefix=f"lo_{session_id}_") as tmp:
        pdf_path = _convert_pptx_to_pdf(pptx_path, Path(tmp))
        image_paths = _render_pdf_to_jpegs(pdf_path, slides_dir)

    if len(image_paths) != len(slide_texts):
        logger.warning(
            "slide count mismatch: images=%d pptx_slides=%d — aligning to min",
            len(image_paths),
            len(slide_texts),
        )
    n = min(len(image_paths), len(slide_texts))
    if n == 0:
        raise ValueError("No slides found in the presentation")

    out: list[PptSlideData] = []
    for i in range(n):
        out.append(
            PptSlideData(
                slide_index=i,
                slide_num=i + 1,
                source_text=slide_texts[i],
                image_path=image_paths[i],
            )
        )
    return out


def build_ppt_segment_dicts(
    slides: list[PptSlideData],
    *,
    session_id: str,
    lang_opt: str,
    subject: str = "",
) -> list[dict]:
    """One segment per slide, with slide_url for UI and vision."""
    from .lang_detect import detect_source_lang, resolve_segment_lang

    segments: list[dict] = []
    for s in slides:
        detected = detect_source_lang(s.source_text)
        sl = resolve_segment_lang(lang_opt, detected)
        text = s.source_text
        fname = s.image_path.name
        segments.append(
            {
                "segment_id": f"s_{s.slide_num:04d}",
                "pages": [s.slide_num],
                "slide_index": s.slide_index,
                "slide_url": f"/api/sessions/{session_id}/slides/{fname}",
                "slide_file": fname,
                "source_text": text,
                "source_lang": sl,
                "chapter_id": "",
                "chapter_title": subject or "Presentation",
                "chapter_index": -1,
            }
        )
    return segments
